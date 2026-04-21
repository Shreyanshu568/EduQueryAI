# CCNA_Rag Chatbot :- 

from langchain_groq import ChatGroq
from langchain_huggingface import HuggingFaceEmbeddings
from langchain_chroma import Chroma
from pptx import Presentation
import os
import re
from dotenv import load_dotenv

load_dotenv()

# Step 1 :- Load PPTX Files :- 

def extract_text_from_pptx(filepaath):
    presentation = Presentation(filepaath)
    text = ""
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text+= paragraph.text + "\n"

    return text 

def load_all_modules(folder):
    all_text = ""
    pptx_files = [f for f in os.listdir(folder) if f.endswith(".pptx")]
    pptx_files.sort()
    for file in pptx_files:
        filepath = os.path.join(folder, file)
        text = extract_text_from_pptx(filepath)
        all_text+= f"\n\n--- {file} ---\n\n" + text


    return all_text, len(pptx_files) 

    # Chunk the text :- 

def chunk_text(text, chunk_size = 500, overlap=50):
     chunks = []
     start = 0
     while start < len(text):
            end=start+chunk_size
            chunk = text[start:end]
            if chunk.strip():
                chunks.append(chunk)

            start = end - overlap
     return chunks 


# Main Execution :- 

if __name__ == "__main__":

    raw_text, total_files = load_all_modules(".")

    chunks = chunk_text(raw_text)

    print(f"Total files loaded: {total_files}")
    print(f"Total characters: {len(raw_text)}")
    print(f"Total chunks created: {len(chunks)}") 


# Step 3 :- Create Vectore Database :- 

def create_vectorstore(chunks):
    embeddings = HuggingFaceEmbeddings(
    model_name="all-MiniLM-L6-v2",
    model_kwargs={"device": "cpu"},
    encode_kwargs={"normalize_embeddings": False}
)

    vectorstore = Chroma.from_texts(
        texts=chunks,
        embedding=embeddings,
        persist_directory="ccna_vectorstore"
    )

    return vectorstore 

def load_vectorstore():
    embeddings = HuggingFaceEmbeddings(
    model_name="all-MiniLM-L6-v2",
    model_kwargs={"device": "cpu"},
    encode_kwargs={"normalize_embeddings": False}
)

    vectorstore = Chroma(
        persist_directory="ccna_vectorstore",
        embedding_function=embeddings
    )

    return vectorstore 

if __name__ == "__main__":

    if not os.path.exists("ccna_vectorstore"):
        print("Creating vector database for first  time...")
        raw_text, total_files = load_all_modules(".")
        chunks = chunk_text(raw_text)
        vectorstore = create_vectorstore(chunks)
        print(f"Done {len(chunks)} chunks stored")

    else:
        print("Loading existing vector database...")
        vectorstore = create_vectorstore(chunks)
        print("Vector database loaded") 


# Step 4 :- Search Function :- 

def search_notes(vectorstore, question, k=15):
    results = vectorstore.similarity_search(question,k=k)

    context = ""
    for i, result in enumerate(results):
        context += f"\n--- Result  {i+1} ---\n"
        context += result.page_content
        context+= "\n"

    return context


# Step 5 :- Generate Answer with LLM


def generate_answer(vectorstore, question):
    
    
    context = search_notes(vectorstore, question)
    
    llm = ChatGroq(
        model="llama-3.3-70b-versatile",
        api_key=os.getenv("GROQ_API_KEY")
    )

    # Answer length preference :- 

    word_count_match = re.search(r'(\d+)\s*words?', question.lower())

    if word_count_match:
        requested_words = word_count_match.group(1)
        word_count = f"Write in exactly {requested_words} words"

    elif any(word in question.lower() for word in ['brief', 'short', 'summarize', 'one line', 'quickly']):
        word_count = "Write in 50-100 words only"

    elif any(word in question.lower() for word in ['detail', 'elaborate', 'comprehensive', 'thorough']):
        word_count = "Write at least 800-1000 words"

    else:
        word_count = "Write at least 400-500 words"
    
    prompt = f"""You are an expert CCNA networking instructor helping students prepare for their cisco certification exam.

    Your job is to give DETAILED, CLEAR and COMPREHENSIVE answers based on the CCNA study materials provided below.

    Instructions:
    - Give thorough explaination with all important details
    - Break down complex topics into simple understandable points
    - Include definitions, examples, and how things work
    - Use numbered points or bullet points for clarity
    - Cover all aspects of the question completely
    - {word_count}
    - Give Complete and thorough explaination so student doesn't need to refer anywhere else 

Context from CCNA notes:
{context}

Student Question: {question}

Comprehensive Answer:"""
    
    response = llm.invoke(prompt)
    
    return response.content

# Test Search :- 

if __name__ == "__main__":
    
    if not os.path.exists("ccna_vectorstore"):
        print("Creating vector database for first time...")
        raw_text, total_files = load_all_modules(".")
        chunks = chunk_text(raw_text)
        vectorstore = create_vectorstore(chunks)
        print(f"Done!! {len(chunks)} chunks stored")
    else:
        vectorstore = load_vectorstore()

    print("\n CCNA Chatbot Ready")
    print("Type 'exit' to quit\n")
    
    while True:
        question = input("You: ")
        
        if question.lower() == "exit":
            print("Goodbye")
            break
        
        if question.strip() == "":
            continue
            
        answer = generate_answer(vectorstore, question)
        print(f"\nChatbot: {answer}\n") 

def generate_quiz_questions(vectorstore, topic):
    context = search_notes(vectorstore, topic)

    llm = ChatGroq(
        model="llama-3.3-70b-versatile",
        api_key=os.getenv("GROQ_API_KEY")
    )

    prompt = f"""You are a CCNA exam question generator.

Generate 5 MCQ questions Exactly like this:
Q1. Question here?
A) Option 1
B) Option 2
C) Option 3
D) Option 4

Answer: B
Explanation: Brief explanation here

Context:
{context}

Topic: {topic}

Generate 5 MCQ questions now:"""
    
    response = llm.invoke(prompt)
    return response.content